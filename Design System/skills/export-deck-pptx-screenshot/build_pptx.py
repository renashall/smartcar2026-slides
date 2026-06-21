#!/usr/bin/env python3
"""Export a rendered HTML slide deck to a 16:9 PPTX of slide screenshots.

Usage:
    python build_pptx.py <deck.html> [out.pptx]

Pipeline (the "screenshot" PPTX):
  1. Render the deck to a clean PDF (reusing a sibling ``<stem>.pdf`` if present,
     e.g. one already made by the export-deck-pdf skill; otherwise a throwaway).
  2. Rasterize each PDF page at 2x to a full-bleed image (pixel-perfect visuals)
     and place it on a 13.333x7.5in (16:9) slide.
  3. Attach the deck's hidden ``<script id="speaker-notes">`` entries as notes.
  4. Add a transparent, clickable hyperlink hotspot over each ``<a href>``,
     positioned from the PDF's own link rect (PDF points map 1:1 to slide points).

Output defaults to ``<deck-stem>.pptx`` beside the HTML.

TRADE-OFF: slide text is a rendered image, not editable text boxes — this buys
exact visual fidelity. Links and speaker notes are preserved.

Requires: python-pptx, pymupdf   (pip install python-pptx pymupdf)
Set the DECK_BROWSER env var to override browser auto-detection for the PDF step.
"""
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
import time
from pathlib import Path

import fitz  # pymupdf
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

BROWSER_CANDIDATES = [
    r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe",
    r"C:\Program Files\Microsoft\Edge\Application\msedge.exe",
    r"C:\Program Files\Google\Chrome\Application\chrome.exe",
    r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe",
    "/Applications/Microsoft Edge.app/Contents/MacOS/Microsoft Edge",
    "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
]
WHICH_NAMES = ["microsoft-edge", "google-chrome", "google-chrome-stable",
               "chromium", "chromium-browser", "chrome", "msedge"]


def find_browser():
    env = os.environ.get("DECK_BROWSER")
    if env and Path(env).exists():
        return env
    for cand in BROWSER_CANDIDATES:
        if Path(cand).exists():
            return cand
    for name in WHICH_NAMES:
        found = shutil.which(name)
        if found:
            return found
    sys.exit("No Edge/Chrome found. Set the DECK_BROWSER env var to a "
             "Chromium-based browser executable.")


def render_pdf(html, pdf, browser=None, attempts=3):
    html, pdf = Path(html), Path(pdf)
    browser = browser or find_browser()
    url = html.resolve().as_uri()
    for _ in range(attempts):
        if pdf.exists():
            pdf.unlink()
        with tempfile.TemporaryDirectory(prefix="deck-pdf-") as profile:
            cmd = [
                browser, "--headless=new", "--disable-gpu", "--no-sandbox",
                "--no-pdf-header-footer", "--run-all-compositor-stages-before-draw",
                "--virtual-time-budget=12000", "--user-data-dir=" + profile,
                "--print-to-pdf=" + str(pdf), url,
            ]
            try:
                subprocess.run(cmd, timeout=90,
                               stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            except subprocess.TimeoutExpired:
                pass
        if pdf.exists() and pdf.stat().st_size > 0:
            return pdf
        time.sleep(1)
    sys.exit("PDF was not produced after %d attempts: %s" % (attempts, pdf))


def speaker_notes(html_text):
    match = re.search(r'id="speaker-notes"[^>]*>\s*(\[.*?\])\s*</script>',
                      html_text, re.S)
    return json.loads(match.group(1)) if match else []


def build(html, pptx):
    html, pptx = Path(html), Path(pptx)
    notes = speaker_notes(html.read_text(encoding="utf-8"))

    sibling_pdf = html.with_suffix(".pdf")
    tmp_pdf_dir = None
    if sibling_pdf.exists() and sibling_pdf.stat().st_size > 0:
        pdf = sibling_pdf
    else:
        tmp_pdf_dir = tempfile.TemporaryDirectory(prefix="deck-pdf-")
        pdf = Path(tmp_pdf_dir.name) / (html.stem + ".pdf")
        render_pdf(html, pdf)

    doc = fitz.open(pdf)
    page_count = doc.page_count
    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 960 pt
    prs.slide_height = Inches(7.5)     # 540 pt
    blank = prs.slide_layouts[6]
    matrix = fitz.Matrix(2, 2)         # 2x for crisp images
    link_count = 0
    with tempfile.TemporaryDirectory(prefix="deck-img-") as imgs:
        for i, page in enumerate(doc):
            png = Path(imgs) / ("p%02d.png" % i)
            page.get_pixmap(matrix=matrix).save(png)
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(str(png), 0, 0,
                                     width=prs.slide_width, height=prs.slide_height)
            # clickable hotspots from the PDF's URI link annotations (pt == pt, 1:1)
            for link in page.get_links():
                uri = link.get("uri")
                if link.get("kind") != fitz.LINK_URI or not uri:
                    continue
                r = link["from"]
                shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                             Pt(r.x0), Pt(r.y0), Pt(r.width), Pt(r.height))
                shp.fill.background()        # transparent
                shp.line.fill.background()   # no outline
                shp.shadow.inherit = False
                shp.click_action.hyperlink.address = uri
                link_count += 1
            if i < len(notes):
                slide.notes_slide.notes_text_frame.text = notes[i]
        doc.close()
        prs.save(pptx)
    if tmp_pdf_dir:
        tmp_pdf_dir.cleanup()
    print("wrote %s  (slides=%d, notes=%d, clickable_links=%d)"
          % (pptx, page_count, len(notes), link_count))


def main():
    if len(sys.argv) < 2:
        sys.exit("Usage: python build_pptx.py <deck.html> [out.pptx]")
    html = Path(sys.argv[1]).resolve()
    if not html.exists():
        sys.exit("HTML not found: %s" % html)
    pptx = Path(sys.argv[2]).resolve() if len(sys.argv) > 2 else html.with_suffix(".pptx")
    build(html, pptx)


if __name__ == "__main__":
    main()
